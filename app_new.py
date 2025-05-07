import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from io import StringIO
from datetime import datetime
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import base64
from pptx.dml.color import RGBColor

# Configurações do Streamlit para permitir upload de arquivos
st.set_option('deprecation.showfileUploaderEncoding', False)

# Configuração da página
st.set_page_config(
    page_title='Dashboard Cubo Itaú',
    layout='wide',
    initial_sidebar_state='collapsed'
)

# Cores do Itaú
CORES_ITAU = {
    'laranja': '#EC7000',
    'azul_escuro': '#003366',
    'azul_claro': '#0057FF',
    'branco': '#FFFFFF',
    'cinza_claro': '#F5F6FA'
}

# =====================
# Função para carregar e pré-processar os dados
# =====================
def carregar_dados():
    st.sidebar.header('Carregar Dados')
    
    # Configuração do file_uploader com todos os tipos de Excel
    uploaded_file = st.sidebar.file_uploader(
        'Faça upload do arquivo Excel',
        type=['xlsx', 'xls', 'xlsm', 'xlsb'],
        accept_multiple_files=False,
        help='Formatos aceitos: .xlsx, .xls, .xlsm, .xlsb'
    )
    
    df = None
    if uploaded_file is not None:
        try:
            # Detecta o tipo de arquivo
            file_type = uploaded_file.name.split('.')[-1].lower()
            
            if file_type == 'xls':
                # Para arquivos Excel 97-2003
                df = pd.read_excel(uploaded_file, engine='xlrd')
            else:
                # Para outros formatos
                try:
                    df = pd.read_excel(uploaded_file, engine='openpyxl')
                except:
                    df = pd.read_excel(uploaded_file, engine='xlrd')
            
            st.sidebar.success('Arquivo carregado com sucesso!')
            
        except Exception as e:
            st.sidebar.error(f'Erro ao ler o arquivo: {str(e)}')
            st.sidebar.info('Dica: Se o arquivo for Excel 97-2003 (.xls), tente salvá-lo como Excel 2007 ou superior (.xlsx)')
            return None
    else:
        st.sidebar.write('Ou cole os dados da planilha (Ctrl+V)')
        clipboard_data = st.sidebar.text_area('Cole aqui os dados copiados da planilha')
        if clipboard_data:
            try:
                df = pd.read_csv(StringIO(clipboard_data), sep='\t')
                st.sidebar.success('Dados colados com sucesso!')
            except Exception as e:
                st.sidebar.error(f'Erro ao ler os dados colados: {str(e)}')
                return None
    
    if df is not None:
        df = preprocessar_dados(df)
    return df

# =====================
# Função de pré-processamento
# =====================
def preprocessar_dados(df):
    # Limpeza da coluna Cliente
    if 'Cliente' in df.columns:
        df['Cliente'] = df['Cliente'].astype(str).str.replace(r'^\d+\s*-\s*', '', regex=True).str.strip()
    
    # Extrair apenas a data da coluna 'Data do Convite' (ex: '30/04/2025 (18:00 às 19:00)' -> '30/04/2025')
    if 'Data do Convite' in df.columns:
        df['Data do Convite'] = df['Data do Convite'].astype(str).str.extract(r'(\d{2}/\d{2}/\d{4})')[0]
    
    # Conversão de datas
    for col in ['Data de Cadastro', 'Data do Convite']:
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], errors='coerce', dayfirst=True)
    
    # Remover linhas com datas inválidas
    df = df.dropna(subset=['Data do Convite'])
    
    # Extrair dia da semana, mês, ano da Data do Convite
    if 'Data do Convite' in df.columns:
        # Garantir que o dia da semana é calculado pela data
        df['Dia da Semana'] = df['Data do Convite'].dt.dayofweek
        nomes_semana = ['Segunda','Terça','Quarta','Quinta','Sexta','Sábado','Domingo']
        df['Dia da Semana Nome'] = df['Dia da Semana'].apply(lambda x: nomes_semana[x] if pd.notnull(x) else '')
        df['Ano'] = df['Data do Convite'].dt.year
        df['Mês'] = df['Data do Convite'].dt.month
        df['Dia'] = df['Data do Convite'].dt.day
    
    return df

# =====================
# Funções para métricas
# =====================
def total_convites(df):
    return len(df)

def anfitrioes_notificados(df):
    return df[df['Anfitrião Notificado'].str.lower() == 'sim'].shape[0]

def anfitrioes_nao_notificados(df):
    return df[df['Anfitrião Notificado'].str.lower() == 'não'].shape[0]

def total_convidados_cubo(df):
    # Cliente segregado: 878 - Cubo
    return df[df['Cliente'].str.lower() == 'cubo'].shape[0]

def total_convidados_residentes(df):
    return len(df) - total_convidados_cubo(df)

def media_convidados_dia_util(df):
    dias_uteis = df[df['Data do Convite'].dt.weekday < 5]['Data do Convite'].dt.date.nunique()
    if dias_uteis == 0:
        return 0
    return int(round(len(df[df['Data do Convite'].dt.weekday < 5]) / dias_uteis, 0))

# =====================
# Funções para gráficos
# =====================
def grafico_top_empresas(df):
    # Remover "Cubo" do top 10
    df_empresas = df[~df['Cliente'].str.lower().str.contains('cubo')]
    top_empresas = df_empresas['Cliente'].value_counts().head(10)
    df_plot = pd.DataFrame({
        'Empresa': top_empresas.index,
        'Convites': top_empresas.values
    })
    fig = px.bar(
        df_plot,
        x='Empresa',
        y='Convites',
        title='Top 10 Empresas com Mais Convites',
        color_discrete_sequence=[CORES_ITAU['azul_escuro']]
    )
    fig.update_traces(text=df_plot['Convites'], textposition='outside')
    fig.update_layout(
        plot_bgcolor=CORES_ITAU['cinza_claro'],
        paper_bgcolor=CORES_ITAU['cinza_claro'],
        title_font_size=22,
        title_font_family='Arial',
        margin=dict(t=60, b=40, l=40, r=40)
    )
    return fig

def grafico_convidados_por_data(df):
    if df.empty:
        return px.bar(title='Sem dados para exibir')
    # Criar range completo de dias do mês selecionado
    data_inicio = df['Data do Convite'].min()
    data_fim = df['Data do Convite'].max()
    if pd.isna(data_inicio) or pd.isna(data_fim):
        return px.bar(title='Datas inválidas')
    mes = data_inicio.month
    ano = data_inicio.year
    dias_no_mes = pd.Period(f'{ano}-{mes:02d}').days_in_month
    dias_mes = pd.date_range(start=f'{ano}-{mes:02d}-01', end=f'{ano}-{mes:02d}-{dias_no_mes}')
    por_data = df.groupby(df['Data do Convite'].dt.date).size()
    por_data = por_data.reindex(dias_mes.date, fill_value=0)
    df_plot = pd.DataFrame({
        'Dia': [str(d.day) for d in por_data.index],
        'Convidados': por_data.values
    })
    fig = px.bar(
        df_plot,
        x='Dia',
        y='Convidados',
        title='Convidados por Dia',
        color_discrete_sequence=[CORES_ITAU['laranja']]
    )
    fig.update_traces(text=df_plot['Convidados'], textposition='outside')
    fig.update_xaxes(tickangle=0, dtick=1, tickmode='array', tickvals=[str(i) for i in range(1, dias_no_mes+1)], ticktext=[str(i) for i in range(1, dias_no_mes+1)])
    fig.update_layout(
        plot_bgcolor=CORES_ITAU['cinza_claro'],
        paper_bgcolor=CORES_ITAU['cinza_claro'],
        title_font_size=22,
        title_font_family='Arial',
        margin=dict(t=60, b=40, l=40, r=40)
    )
    return fig

def grafico_convidados_por_dia_semana(df):
    nomes_semana = ['Segunda','Terça','Quarta','Quinta','Sexta','Sábado','Domingo']
    por_dia = df['Dia da Semana Nome'].value_counts().reindex(nomes_semana, fill_value=0)
    df_plot = pd.DataFrame({
        'Dia da Semana': nomes_semana,
        'Convidados': por_dia.values
    })
    fig = px.bar(
        df_plot,
        x='Dia da Semana',
        y='Convidados',
        title='Convidados por Dia da Semana',
        color_discrete_sequence=[CORES_ITAU['azul_escuro']]
    )
    fig.update_traces(text=df_plot['Convidados'], textposition='outside')
    fig.update_layout(
        plot_bgcolor=CORES_ITAU['cinza_claro'],
        paper_bgcolor=CORES_ITAU['cinza_claro'],
        title_font_size=22,
        title_font_family='Arial',
        margin=dict(t=60, b=40, l=40, r=40)
    )
    return fig

# =====================
# Visitantes Frequentes por Empresa (>4 visitas no mês)
# =====================
def visitantes_frequentes(df):
    tabela = []
    for empresa, grupo in df.groupby('Cliente'):
        visitantes = grupo.groupby('E-mail').size()
        frequentes = visitantes[visitantes > 4]
        if not frequentes.empty:
            for email, qtd in frequentes.items():
                tabela.append({'Empresa': empresa, 'E-mail': email, 'Visitas': qtd})
    df_tabela = pd.DataFrame(tabela)
    if not df_tabela.empty:
        df_tabela = df_tabela.sort_values('Visitas', ascending=False)
    return df_tabela

def consolidado_frequentes(df):
    tabela = visitantes_frequentes(df)
    if tabela.empty:
        return pd.DataFrame(columns=['Quantidade de Empresas', 'Ocorrências'])
    # Conta quantas empresas tiveram X visitantes frequentes
    ocorrencias = tabela.groupby('Empresa').size().value_counts().sort_index()
    return pd.DataFrame({
        'Ocorrências': ocorrencias.index,
        'Quantidade de Empresas': ocorrencias.values
    })

def consolidado_frequentes_grafico(df):
    tabela = visitantes_frequentes(df)
    if tabela.empty:
        return None
    ocorrencias = tabela.groupby('Empresa').size().value_counts().sort_index()
    fig = go.Figure(go.Bar(
        x=ocorrencias.values,
        y=[f"{i} visitantes" for i in ocorrencias.index],
        orientation='h',
        marker_color=CORES_ITAU['azul_escuro'],
        text=ocorrencias.values,
        textposition='outside'
    ))
    fig.update_layout(
        title='Empresas por quantidade de visitantes frequentes',
        xaxis_title='Quantidade de Empresas',
        yaxis_title='',
        plot_bgcolor=CORES_ITAU['cinza_claro'],
        paper_bgcolor=CORES_ITAU['cinza_claro'],
        height=300
    )
    return fig

def painel_empresas_frequentes(df):
    tabela = visitantes_frequentes(df)
    if tabela.empty:
        return ''
    ocorrencias = tabela.groupby('Empresa').size()
    painel = ''
    for qtd in sorted(ocorrencias.unique()):
        empresas = ocorrencias[ocorrencias == qtd].index.tolist()
        empresas_str = ', '.join(empresas)
        painel += f'<div style="margin-bottom:12px;"><b>{qtd} visitantes frequentes</b><br><span style="color:#003366">{empresas_str}</span></div>'
    return painel

def gerar_pptx(df, df_filtro):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    left = Inches(0.2)
    top = Inches(0.2)
    width = Inches(2.2)
    height = Inches(0.7)
    azul_escuro = (0, 51, 102)  # RGB do Itaú
    # Adiciona cards
    for i, (label, value) in enumerate([
        ("Total de Convites", total_convites(df)),
        ("Anfitriões Notificados", anfitrioes_notificados(df)),
        ("Não Notificados", anfitrioes_nao_notificados(df)),
        ("Convidados Cubo", total_convidados_cubo(df)),
        ("Média por Dia Útil", media_convidados_dia_util(df)),
    ]):
        txBox = slide.shapes.add_textbox(left + Inches(i*2.3), top, width, height)
        tf = txBox.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*azul_escuro)
        p2 = tf.add_paragraph()
        p2.text = str(value)
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(*azul_escuro)
        p2.alignment = PP_ALIGN.CENTER
    # Adiciona título
    titleBox = slide.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(10), Inches(0.7))
    titleBox.text = 'Dashboard de Visitas - Cubo Itaú'
    # Exporta slide para bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

def main():
    st.markdown(f"""
        <style>
        .modern-card {{
            background: {CORES_ITAU['cinza_claro']};
            border-radius: 18px;
            padding: 24px 16px 16px 16px;
            margin-bottom: 12px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.04);
            text-align: center;
            min-width: 250px;
            max-width: 250px;
            height: 140px;
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
        }}
        .modern-card .big-number {{
            font-size: 2.5em;
            font-weight: bold;
            color: {CORES_ITAU['azul_escuro']};
        }}
        .modern-card .card-label {{
            font-size: 1.2em;
            color: {CORES_ITAU['laranja']};
            font-weight: 600;
        }}
        .main-title {{
            font-size: 2.5em;
            font-weight: bold;
            color: {CORES_ITAU['azul_escuro']};
            margin-bottom: 1em;
        }}
        .frequent-visitors-table {{
            background: {CORES_ITAU['cinza_claro']};
            border-radius: 18px;
            padding: 24px;
            margin: 0;
            box-shadow: 0 2px 8px rgba(0,0,0,0.04);
        }}
        </style>
    """, unsafe_allow_html=True)
    
    st.markdown('<div class="main-title">Dashboard de Visitas - Cubo Itaú</div>', unsafe_allow_html=True)
    
    df = carregar_dados()
    if df is None:
        st.info('Por favor, carregue um arquivo Excel ou cole os dados para iniciar a análise.')
        return

    # Cards em linha horizontal usando st.columns, igualmente espaçados
    col1, col2, col3, col4, col5 = st.columns(5)
    with col1:
        st.markdown(f'<div class="modern-card"><div class="card-label">Total de Convites</div><div class="big-number">{total_convites(df)}</div></div>', unsafe_allow_html=True)
    with col2:
        st.markdown(f'<div class="modern-card"><div class="card-label">Anfitriões Notificados</div><div class="big-number">{anfitrioes_notificados(df)}</div></div>', unsafe_allow_html=True)
    with col3:
        st.markdown(f'<div class="modern-card"><div class="card-label">Não Notificados</div><div class="big-number">{anfitrioes_nao_notificados(df)}</div></div>', unsafe_allow_html=True)
    with col4:
        st.markdown(f'<div class="modern-card"><div class="card-label">Convidados Cubo</div><div class="big-number">{total_convidados_cubo(df)}</div></div>', unsafe_allow_html=True)
    with col5:
        st.markdown(f'<div class="modern-card"><div class="card-label">Média por Dia Útil</div><div class="big-number">{media_convidados_dia_util(df)}</div></div>', unsafe_allow_html=True)

    # Filtro de período
    st.sidebar.header('Filtro de Período')
    if df.empty or 'Ano' not in df.columns or 'Mês' not in df.columns:
        st.error('Dados inválidos ou incompletos. Verifique se o arquivo contém as colunas necessárias.')
        return

    anos = sorted(df['Ano'].dropna().unique(), reverse=True)
    meses = sorted(df['Mês'].dropna().unique())
    if not anos or not meses:
        st.error('Não há dados de período disponíveis.')
        return

    ano_sel = st.sidebar.selectbox('Ano', anos)
    mes_sel = st.sidebar.selectbox('Mês', meses)
    df_filtro = df[(df['Ano'] == ano_sel) & (df['Mês'] == mes_sel)]
    if df_filtro.empty:
        st.warning('Não há dados para o período selecionado.')
        return

    # Primeira linha de gráficos (2 colunas)
    col1, col2 = st.columns(2)
    with col1:
        st.plotly_chart(grafico_top_empresas(df_filtro), use_container_width=True)
    with col2:
        st.plotly_chart(grafico_convidados_por_data(df_filtro), use_container_width=True)

    # Segunda linha de gráficos (2 colunas)
    col1, col2 = st.columns(2)
    with col1:
        st.plotly_chart(grafico_convidados_por_dia_semana(df_filtro), use_container_width=True)
    with col2:
        st.markdown('<div class="frequent-visitors-table">', unsafe_allow_html=True)
        st.subheader('Visitantes Frequentes por Empresa (>4 visitas no mês)')
        tabela_frequentes = visitantes_frequentes(df_filtro)
        if not tabela_frequentes.empty:
            st.dataframe(tabela_frequentes, height=420)
        st.markdown('</div>', unsafe_allow_html=True)

    # Terceira seção (consolidado e painel)
    st.markdown('---')
    col1, col2 = st.columns(2)
    with col1:
        st.subheader('Consolidado de Empresas com Visitantes Frequentes')
        st.dataframe(consolidado_frequentes(df_filtro), height=200)
        fig_consolidado = consolidado_frequentes_grafico(df_filtro)
        if fig_consolidado:
            st.plotly_chart(fig_consolidado, use_container_width=True)
    with col2:
        st.subheader('Painel de Empresas com Visitantes Frequentes')
        st.markdown(painel_empresas_frequentes(df_filtro), unsafe_allow_html=True)

    # Botão para download em PPTX
    if st.button('Baixar visualização em PPTX'):
        pptx_bytes = gerar_pptx(df, df_filtro)
        st.download_button('Clique aqui para baixar o PPTX', pptx_bytes, file_name='dashboard_cubo.pptx')

if __name__ == '__main__':
    main() 